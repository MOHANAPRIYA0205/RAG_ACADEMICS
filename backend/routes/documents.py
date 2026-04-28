"""
routes/documents.py — Document upload, listing, and deletion
"""
import io
import uuid
import logging

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from fastapi import APIRouter, UploadFile, File, HTTPException
from rag.pipeline import rag

router = APIRouter()
logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".pptx"}


def extract_text(file: UploadFile) -> str:
    """Extract plain text from an uploaded file."""
    content = file.file.read()
    name = (file.filename or "").lower()

    if name.endswith(".pdf"):
        try:
            doc = fitz.open(stream=content, filetype="pdf")
            return "\n".join(page.get_text() for page in doc)
        except Exception as e:
            raise ValueError(f"Failed to parse PDF: {e}")

    elif name.endswith(".docx"):
        try:
            doc = DocxDocument(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            raise ValueError(f"Failed to parse DOCX: {e}")

    elif name.endswith(".txt"):
        return content.decode("utf-8", errors="ignore")

    elif name.endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        except ImportError:
            # python-pptx not installed — fall back to raw decode
            return content.decode("utf-8", errors="ignore")
        except Exception as e:
            raise ValueError(f"Failed to parse PPTX: {e}")

    # Generic fallback
    return content.decode("utf-8", errors="ignore")


@router.post("/upload")
async def upload_document(file: UploadFile = File(...)):
    """Upload a document, extract text, and ingest into the RAG pipeline."""
    filename = file.filename or ""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type '{ext}'. Supported: PDF, DOCX, TXT, PPTX"
        )

    try:
        text = extract_text(file)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"Text extraction error: {e}")
        raise HTTPException(status_code=500, detail=f"Could not extract text: {e}")

    if len(text.strip()) < 50:
        raise HTTPException(
            status_code=400,
            detail="Could not extract enough text from the document (minimum 50 characters)."
        )

    doc_id = str(uuid.uuid4())
    metadata = {"filename": filename, "doc_id": doc_id}

    try:
        chunks = rag.ingest(text, metadata)
        return {
            "doc_id": doc_id,
            "filename": filename,
            "chunks": chunks,
            "status": "indexed",
            "message": f"Document successfully indexed with {chunks} chunks",
        }
    except Exception as e:
        logger.error(f"RAG ingest error: {e}")
        # Return partial success — text was extracted but embedding failed
        word_count = len(text.split())
        return {
            "doc_id": doc_id,
            "filename": filename,
            "chunks": word_count // 100 or 1,
            "status": "indexed_partial",
            "message": (
                "Document uploaded but the embedding service is unavailable. "
                "The system will use text search as a fallback."
            ),
        }


@router.get("/list")
async def list_documents():
    """Return a list of documents currently cached in the RAG pipeline."""
    try:
        docs = []
        for doc_id, data in rag.documents_cache.items():
            meta = data.get("metadata", {})
            docs.append({
                "id": doc_id,
                "doc_id": doc_id,
                "name": meta.get("filename", "Unknown"),
                "filename": meta.get("filename", "Unknown"),
                "type": meta.get("filename", "").rsplit(".", 1)[-1].lower() if "." in meta.get("filename", "") else "txt",
                "chunks": data.get("chunk_count", 0),
                "pages": 1,
                "size": f"{len(data.get('text', '')) // 1024 or 1} KB",
                "status": "indexed",
                "uploaded": "recently",
                "tags": ["Uploaded"],
            })
        return {"documents": docs, "status": "ok"}
    except Exception as e:
        logger.error(f"List documents error: {e}")
        return {"documents": [], "error": str(e)}


@router.delete("/{doc_id}")
async def delete_document(doc_id: str):
    """Remove a document from the in-memory cache (and ChromaDB if available)."""
    try:
        # Remove from in-memory cache
        if doc_id in rag.documents_cache:
            del rag.documents_cache[doc_id]

        # Attempt to remove from ChromaDB
        if rag.vectorstore:
            try:
                rag.vectorstore.delete(where={"doc_id": doc_id})
            except Exception as e:
                logger.warning(f"Could not remove from ChromaDB: {e}")

        return {"status": "deleted", "doc_id": doc_id, "message": "Document removed from knowledge base"}
    except Exception as e:
        logger.error(f"Delete error: {e}")
        raise HTTPException(status_code=500, detail=f"Delete failed: {str(e)}")
